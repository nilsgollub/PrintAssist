"""
Office document manipulation tools.
Uses python-docx (Word), openpyxl (Excel), python-pptx (PowerPoint).

Routing hints (natural language triggers for Claude Code):
  Word (.docx):
    Text ersetzen / suchen & ersetzen → replace_text_docx()
    Schriftart / Größe ändern         → format_paragraph_docx()
    Absatz einfügen                   → insert_paragraph_docx()

  Excel (.xlsx):
    Zellwert lesen / schreiben        → read_cell_xlsx() / write_cell_xlsx()
    Zeile einfügen                    → insert_row_xlsx()
    Suchen & Ersetzen                 → replace_value_xlsx()

  PowerPoint (.pptx):
    Text in Folie ersetzen            → replace_text_pptx()
    Folie einfügen / löschen         → add_slide_pptx() / delete_slide_pptx()
    Hintergrundfarbe setzen           → set_slide_background_pptx()

NOTE: Converting Office to PDF for printing is done via LibreOffice headless
      (see scripts/print.py). Call print_file() from print.py directly.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import docx
from docx.shared import Pt, RGBColor
import openpyxl
from openpyxl.styles import Font as XlFont
from pptx import Presentation
from pptx.util import Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB


# ===========================================================================
# Word (.docx)
# ===========================================================================

def replace_text_docx(
    input_path: str | Path,
    output_path: str | Path,
    replacements: dict[str, str],
    case_sensitive: bool = True,
) -> Path:
    """
    Find and replace text throughout a Word document (paragraphs + tables).

    Natural language triggers:
      "Text ersetzen in Word", "find and replace", "Platzhalter ausfüllen",
      "{{Name}} durch echten Namen ersetzen", "suchen und ersetzen"

    Args:
        input_path:     Source .docx file.
        output_path:    Where to save the modified file.
        replacements:   Dict of {old_text: new_text}.
        case_sensitive: Whether the search is case-sensitive (default: True).

    Returns:
        Path to the saved file.
    """
    doc = docx.Document(str(input_path))

    def _replace_in_run(run: docx.text.run.Run) -> None:
        for old, new in replacements.items():
            text = run.text
            if not case_sensitive:
                import re
                run.text = re.sub(re.escape(old), new, text, flags=re.IGNORECASE)
            else:
                run.text = text.replace(old, new)

    for para in doc.paragraphs:
        for run in para.runs:
            _replace_in_run(run)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        _replace_in_run(run)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return out


def format_paragraph_docx(
    input_path: str | Path,
    output_path: str | Path,
    paragraph_index: int,
    font_name: str | None = None,
    font_size_pt: float | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    color_rgb: tuple[int, int, int] | None = None,
) -> Path:
    """
    Change formatting of a specific paragraph (0-based index).

    Natural language triggers:
      "Schrift ändern", "fett machen", "Absatz formatieren",
      "Überschrift Schriftgröße anpassen", "Farbe des Textes ändern"

    Args:
        input_path:      Source .docx file.
        output_path:     Where to save the result.
        paragraph_index: 0-based index of the paragraph to format.
        font_name:       Font family name, e.g. "Arial".
        font_size_pt:    Font size in points.
        bold:            True/False or None (unchanged).
        italic:          True/False or None (unchanged).
        color_rgb:       Tuple (R, G, B), each 0–255.

    Returns:
        Path to the saved file.
    """
    doc = docx.Document(str(input_path))
    paras = doc.paragraphs
    if paragraph_index >= len(paras):
        raise IndexError(f"Absatzindex {paragraph_index} existiert nicht (nur {len(paras)} Absätze).")

    para = paras[paragraph_index]
    for run in para.runs:
        if font_name is not None:
            run.font.name = font_name
        if font_size_pt is not None:
            run.font.size = Pt(font_size_pt)
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if color_rgb is not None:
            run.font.color.rgb = RGBColor(*color_rgb)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return out


def insert_paragraph_docx(
    input_path: str | Path,
    output_path: str | Path,
    text: str,
    after_index: int | None = None,
    style: str = "Normal",
) -> Path:
    """
    Insert a new paragraph into a Word document.

    Natural language triggers:
      "Absatz einfügen", "Text hinzufügen", "neuen Abschnitt einfügen",
      "Zeile ergänzen", "append paragraph", "insert text"

    Args:
        input_path:   Source .docx file.
        output_path:  Where to save the result.
        text:         Text content of the new paragraph.
        after_index:  Insert after this 0-based paragraph index. None = append.
        style:        Word style name (default: "Normal").

    Returns:
        Path to the saved file.
    """
    doc = docx.Document(str(input_path))

    if after_index is None or after_index >= len(doc.paragraphs) - 1:
        doc.add_paragraph(text, style=style)
    else:
        ref_para = doc.paragraphs[after_index]
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        new_para = OxmlElement("w:p")
        ref_para._element.addnext(new_para)
        new_doc_para = docx.text.paragraph.Paragraph(new_para, doc)
        run = new_doc_para.add_run(text)
        new_doc_para.style = doc.styles[style]

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return out


# ===========================================================================
# Excel (.xlsx)
# ===========================================================================

def read_cell_xlsx(
    input_path: str | Path,
    sheet_name: str,
    cell: str,
) -> Any:
    """
    Read a single cell value from an Excel workbook.

    Natural language triggers:
      "Zellwert lesen", "was steht in B3", "read cell", "Wert aus Excel holen"

    Args:
        input_path: Source .xlsx file.
        sheet_name: Name of the worksheet.
        cell:       Cell address, e.g. "B3".

    Returns:
        The cell value (str, int, float, datetime, or None).
    """
    wb = openpyxl.load_workbook(str(input_path), read_only=True, data_only=True)
    ws = wb[sheet_name]
    return ws[cell].value


def write_cell_xlsx(
    input_path: str | Path,
    output_path: str | Path,
    sheet_name: str,
    cell: str,
    value: Any,
    bold: bool = False,
    font_size_pt: float | None = None,
) -> Path:
    """
    Write a value into a specific cell of an Excel workbook.

    Natural language triggers:
      "Zellwert schreiben", "Wert eintragen", "update cell", "Zelle befüllen",
      "in Excel ändern", "Spalte aktualisieren"

    Args:
        input_path:   Source .xlsx file.
        output_path:  Where to save the result.
        sheet_name:   Name of the worksheet.
        cell:         Cell address, e.g. "B3".
        value:        New value for the cell.
        bold:         Make cell font bold.
        font_size_pt: Font size in points (None = unchanged).

    Returns:
        Path to the saved file.
    """
    wb = openpyxl.load_workbook(str(input_path))
    ws = wb[sheet_name]
    ws[cell] = value

    if bold or font_size_pt is not None:
        current = ws[cell].font
        ws[cell].font = XlFont(
            bold=bold if bold else current.bold,
            size=font_size_pt if font_size_pt is not None else current.size,
        )

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out))
    return out


def replace_value_xlsx(
    input_path: str | Path,
    output_path: str | Path,
    sheet_name: str,
    old_value: Any,
    new_value: Any,
) -> Path:
    """
    Replace all occurrences of a value in a worksheet.

    Natural language triggers:
      "Excel suchen und ersetzen", "Wert ersetzen in Excel",
      "find and replace Excel", "alle Vorkommen ändern"

    Args:
        input_path:  Source .xlsx file.
        output_path: Where to save the result.
        sheet_name:  Name of the worksheet to search.
        old_value:   Value to find.
        new_value:   Replacement value.

    Returns:
        Path to the saved file.
    """
    wb = openpyxl.load_workbook(str(input_path))
    ws = wb[sheet_name]

    for row in ws.iter_rows():
        for cell in row:
            if cell.value == old_value:
                cell.value = new_value

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out))
    return out


def insert_row_xlsx(
    input_path: str | Path,
    output_path: str | Path,
    sheet_name: str,
    row_index: int,
    values: list[Any],
) -> Path:
    """
    Insert a new row at a given position (1-based) and fill it with values.

    Natural language triggers:
      "Zeile einfügen", "neue Zeile hinzufügen", "insert row",
      "Datenzeile ergänzen", "Tabellenzeile hinzufügen"

    Args:
        input_path:  Source .xlsx file.
        output_path: Where to save the result.
        sheet_name:  Name of the worksheet.
        row_index:   1-based row index at which to insert.
        values:      List of cell values for the new row (left to right).

    Returns:
        Path to the saved file.
    """
    wb = openpyxl.load_workbook(str(input_path))
    ws = wb[sheet_name]
    ws.insert_rows(row_index)
    for col_idx, val in enumerate(values, start=1):
        ws.cell(row=row_index, column=col_idx, value=val)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out))
    return out


# ===========================================================================
# PowerPoint (.pptx)
# ===========================================================================

def replace_text_pptx(
    input_path: str | Path,
    output_path: str | Path,
    replacements: dict[str, str],
) -> Path:
    """
    Find and replace text across all slides of a PowerPoint presentation.

    Natural language triggers:
      "Text in PowerPoint ersetzen", "Platzhalter in Folien ausfüllen",
      "replace text in slides", "Folientext ändern", "Präsentation anpassen"

    Args:
        input_path:   Source .pptx file.
        output_path:  Where to save the result.
        replacements: Dict of {old_text: new_text}.

    Returns:
        Path to the saved file.
    """
    prs = Presentation(str(input_path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        run.text = run.text.replace(old, new)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def add_slide_pptx(
    input_path: str | Path,
    output_path: str | Path,
    layout_index: int = 1,
    title: str = "",
    body: str = "",
) -> Path:
    """
    Append a new slide with optional title and body text.

    Natural language triggers:
      "Folie hinzufügen", "neue Slide einfügen", "add slide",
      "Seite ergänzen in PowerPoint", "Folie am Ende einfügen"

    Args:
        input_path:    Source .pptx file.
        output_path:   Where to save the result.
        layout_index:  Slide layout index (0-based, default: 1 = title+content).
        title:         Text for the title placeholder.
        body:          Text for the body/content placeholder.

    Returns:
        Path to the saved file.
    """
    prs = Presentation(str(input_path))
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    if title and slide.shapes.title:
        slide.shapes.title.text = title

    if body:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = body
                break

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def delete_slide_pptx(
    input_path: str | Path,
    output_path: str | Path,
    slide_index: int,
) -> Path:
    """
    Delete a slide by its 0-based index.

    Natural language triggers:
      "Folie löschen", "Slide entfernen", "delete slide",
      "Seite X aus Präsentation entfernen"

    Args:
        input_path:   Source .pptx file.
        output_path:  Where to save the result.
        slide_index:  0-based index of the slide to remove.

    Returns:
        Path to the saved file.
    """
    prs = Presentation(str(input_path))
    slides = prs.slides

    if slide_index >= len(slides):
        raise IndexError(f"Folie {slide_index} existiert nicht ({len(slides)} Folien total).")

    xml_slides = slides._sldIdLst
    xml_slides.remove(xml_slides[slide_index])

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def set_slide_background_pptx(
    input_path: str | Path,
    output_path: str | Path,
    slide_index: int,
    color_rgb: tuple[int, int, int],
) -> Path:
    """
    Set a solid background color for a specific slide.

    Natural language triggers:
      "Folienhintergrund ändern", "Hintergrundfarbe setzen",
      "slide background color", "Hintergrund einfärben"

    Args:
        input_path:   Source .pptx file.
        output_path:  Where to save the result.
        slide_index:  0-based slide index.
        color_rgb:    Background color as (R, G, B), each 0–255.

    Returns:
        Path to the saved file.
    """
    prs = Presentation(str(input_path))
    slide = prs.slides[slide_index]
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptRGB(*color_rgb)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
